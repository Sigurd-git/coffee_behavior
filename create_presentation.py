#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
import glob
from pptx import Presentation
from pptx.util import Inches


def create_slide_with_image(prs, image_path, title=None):
    """
    Create a new slide with an image and optional title.

    Args:
        prs: Presentation object
        image_path: Path to the image file
        title: Optional title for the slide

    Returns:
        The created slide
    """
    # Use a slide layout with a title and content
    slide_layout = prs.slide_layouts[5]  # Layout: Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Add title if provided
    if title:
        title_shape = slide.shapes.title
        title_shape.text = title

    # Add the image - reducing size to ensure it fits within bounds
    left = Inches(2)  # Increased left margin for better centering
    top = Inches(1.8)  # Moved down slightly to accommodate title
    width = Inches(6)  # Reduced width for smaller image size
    pic = slide.shapes.add_picture(image_path, left, top, width=width)

    # Add filename at the bottom for reference
    filename = os.path.basename(image_path)
    left = Inches(2)  # Match image left margin
    top = Inches(6.5)  # Moved up further to ensure it's within bounds
    width = Inches(6)  # Match image width
    height = Inches(0.4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = filename

    return slide


def main():
    # Define the output directory and PowerPoint file name
    output_dir = "output"
    ppt_filename = "experiment_results.pptx"

    # Create a new presentation
    prs = Presentation()

    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Get all PNG files in the output directory
    image_files = glob.glob(os.path.join(output_dir, "*.png"))

    # Sort the files
    image_files.sort()

    # Define task types and conditions for organizing slides
    task_types = ["bart", "stroop", "nback", "emotion"]
    conditions = ["control", "exp"]
    metrics = [
        "accuracy",
        "rt",
        "mean_per_balloon",
        "total_earnings",
        "mean_arousal",
        "mean_valence",
    ]

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Experiment Results"
    subtitle.text = "Visual Analysis of Cognitive Tasks"

    # Add section slides and content slides organized by task type and condition
    for task in task_types:
        # Create a section slide for each task
        section_slide_layout = prs.slide_layouts[1]  # Layout with title and content
        section_slide = prs.slides.add_slide(section_slide_layout)
        section_title = section_slide.shapes.title
        section_title.text = f"{task.upper()} Task Results"

        # Find images related to this task
        task_images = [
            img for img in image_files if task in os.path.basename(img).lower()
        ]

        for metric in metrics:
            # Skip metrics not applicable to this task
            if (
                task != "bart"
                and (metric == "mean_per_balloon" or metric == "total_earnings")
            ) or (
                task != "emotion"
                and (metric == "mean_arousal" or metric == "mean_valence")
            ):
                continue

            for condition in conditions:
                # Find images matching this task, metric, and condition
                matching_images = [
                    img
                    for img in task_images
                    if metric in os.path.basename(img).lower()
                    and condition in os.path.basename(img).lower()
                ]

                for image_path in matching_images:
                    # Generate a descriptive title from the filename
                    filename = os.path.basename(image_path)
                    title_parts = filename.replace(".png", "").split("_")

                    if "comparison" in filename:
                        title = f"{task.upper()} Task: {metric.replace('_', ' ').title()} Comparison ({condition.title()})"
                    else:
                        # For other metrics like mean_per_balloon
                        metric_desc = " ".join(
                            [
                                part.title()
                                for part in title_parts
                                if part not in [task, condition]
                            ]
                        )
                        title = (
                            f"{task.upper()} Task: {metric_desc} ({condition.title()})"
                        )

                    create_slide_with_image(prs, image_path, title)

    # Save the PowerPoint file
    prs.save(ppt_filename)
    print(f"Presentation created: {ppt_filename}")


if __name__ == "__main__":
    main()
